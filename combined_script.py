import os
import json
from pdfminer.high_level import extract_text
import docx
from pptx import Presentation
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
import torch
from transformers import BertTokenizer, BertModel
import arabic_reshaper
from bidi.algorithm import get_display

nltk.download('punkt')
nltk.download('stopwords')

# Directories and file paths
FILE_DIR = 'moodle_files'
TEXT_DIR = 'text_files'
JSON_FILE = 'text_data.json'
EMBEDDINGS_FILE = 'embeddings.json'

# Initialize BERT tokenizer and model for Arabic
tokenizer = BertTokenizer.from_pretrained('asafaya/bert-base-arabic')
model = BertModel.from_pretrained('asafaya/bert-base-arabic')

def get_tokenizer_and_model():
    return tokenizer, model

def convert_pdf_to_text(pdf_path):
    try:
        return extract_text(pdf_path)
    except Exception as e:
        print(f"Error converting {pdf_path}: {e}")
        return None

def convert_docx_to_text(docx_path):
    try:
        doc = docx.Document(docx_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Error converting {docx_path}: {e}")
        return None

def convert_ppt_to_text(ppt_path):
    try:
        prs = Presentation(ppt_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Error converting {ppt_path}: {e}")
        return None

def preprocess_text(text):
    reshaped_text = arabic_reshaper.reshape(text)
    bidi_text = get_display(reshaped_text)
    tokens = word_tokenize(bidi_text)
    stop_words = set(stopwords.words('arabic'))
    filtered_tokens = [word for word in tokens if word.lower() not in stop_words]
    return ' '.join(filtered_tokens)

def save_text_data(filename, text):
    if not os.path.exists(TEXT_DIR):
        os.makedirs(TEXT_DIR)
    with open(os.path.join(TEXT_DIR, filename), 'w', encoding='utf-8') as f:
        f.write(text)

def load_text_files(text_dir):
    text_data = {}
    for text_file in os.listdir(text_dir):
        if text_file.endswith('.txt'):
            file_path = os.path.join(text_dir, text_file)
            with open(file_path, 'r', encoding='utf-8') as f:
                text_data[text_file] = f.read()
    return text_data

def save_to_json(data, json_file):
    with open(json_file, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=4)

def encode_texts(texts):
    inputs = tokenizer(texts, return_tensors='pt', padding=True, truncation=True, max_length=512)
    with torch.no_grad():
        outputs = model(**inputs)
        embeddings = outputs.last_hidden_state.mean(dim=1)  # Average pooling
    return embeddings

def main():
    # Step 1: Convert files to text
    for file in os.listdir(FILE_DIR):
        file_path = os.path.join(FILE_DIR, file)
        text = None
        if file.endswith('.pdf'):
            print(f'Converting {file} to text...')
            text = convert_pdf_to_text(file_path)
        elif file.endswith('.docx'):
            print(f'Converting {file} to text...')
            text = convert_docx_to_text(file_path)
        elif file.endswith('.ppt') or file.endswith('.pptx'):
            print(f'Converting {file} to text...')
            text = convert_ppt_to_text(file_path)
        
        if text:
            print(f"Extracted text from {file}:\n{text[:500]}...")  # Print first 500 characters of extracted text for debugging
            preprocessed_text = preprocess_text(text)
            save_text_data(file.replace(file.split('.')[-1], 'txt'), preprocessed_text)
            print(f'Saved text for {file}.')

    # Step 2: Load text files and save to JSON
    text_data = load_text_files(TEXT_DIR)
    save_to_json(text_data, JSON_FILE)
    print(f'Stored text data in {JSON_FILE}')

    # Step 3: Encode texts and save embeddings
    texts = list(text_data.values())
    embeddings = encode_texts(texts)
    embeddings_dict = {file: embedding.tolist() for file, embedding in zip(text_data.keys(), embeddings)}

    with open(EMBEDDINGS_FILE, 'w', encoding='utf-8') as f:
        json.dump(embeddings_dict, f, ensure_ascii=False, indent=4)

    print(f'Embeddings stored in {EMBEDDINGS_FILE}')

if __name__ == '__main__':
    main()